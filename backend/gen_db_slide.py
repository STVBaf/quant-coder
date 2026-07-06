"""Generate a single DB-design slide (docs/数据库设计.pptx), same dark theme as
the main deck. Standalone so it can be inserted into the deck manually.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG        = RGBColor(0x0B, 0x0E, 0x11)
PANEL     = RGBColor(0x18, 0x1A, 0x20)
LINE      = RGBColor(0x2B, 0x31, 0x39)
FG        = RGBColor(0xEA, 0xEC, 0xEF)
FG_STRONG = RGBColor(0xFA, 0xFA, 0xFA)
FG_DIM    = RGBColor(0x84, 0x8E, 0x9C)
UP        = RGBColor(0xF6, 0x46, 0x5D)
DOWN      = RGBColor(0x0E, 0xCB, 0x81)
BRAND     = RGBColor(0x00, 0x7A, 0xFF)
AI        = RGBColor(0x8A, 0x2B, 0xE2)
FONT      = "Microsoft YaHei"
FONT_MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _ns(sp):
    sp.shadow.inherit = False


def rect(s, x, y, w, h, fill=PANEL, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    _ns(sp)
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        for (t, size, color, bold, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = fnt
    return tb


def node(s, x, y, w, name, desc, color, h=Inches(0.92)):
    """Simple E-R entity node: table name + one-line description (no fields)."""
    rect(s, x, y, w, h, fill=PANEL, line=color, line_w=1.75)
    rect(s, x, y, Inches(0.08), h, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)
    txt(s, x + Inches(0.24), y + Inches(0.12), w - Inches(0.35), Inches(0.42),
        [[(name, 15, FG_STRONG, True, FONT_MONO)]])
    txt(s, x + Inches(0.24), y + Inches(0.52), w - Inches(0.35), Inches(0.35),
        [[(desc, 11, FG_DIM, False, FONT)]])
    return (x, y, w, h)


def connect(s, a, b, label, side="right"):
    """Draw a 1:N connector line between two nodes with a label near it."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    if side == "right":  # a right edge -> b left edge
        x1, y1 = ax + aw, ay + Emu(ah.emu // 2)
        x2, y2 = bx, by + Emu(bh.emu // 2)
    else:  # vertical: a bottom -> b top
        x1, y1 = ax + Emu(aw.emu // 2), ay + ah
        x2, y2 = bx + Emu(bw.emu // 2), by
    ln = s.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    ln.line.color.rgb = LINE
    ln.line.width = Pt(1.5)
    _ns(ln)
    # label at midpoint
    mx, my = Emu((x1 + x2) // 2), Emu((y1 + y2) // 2)
    txt(s, mx - Inches(0.35), my - Inches(0.22), Inches(0.9), Inches(0.3),
        [[(label, 11, BRAND, True, FONT_MONO)]], align=PP_ALIGN.CENTER)


# ── Slide ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background(); _ns(bg)

# chrome + title
txt(s, Inches(0.6), Inches(0.34), Inches(6), Inches(0.4),
    [[("QuantDesk", 12, BRAND, True, FONT), ("  ·  数据库设计", 12, FG_DIM, False, FONT)]])
txt(s, Inches(11.7), Inches(6.95), Inches(1.1), Inches(0.4),
    [[("DB", 11, FG_DIM, False, FONT_MONO)]], align=PP_ALIGN.RIGHT)
rect(s, Inches(0.6), Inches(1.0), Inches(0.08), Inches(0.62), fill=BRAND, line=None, shape=MSO_SHAPE.RECTANGLE)
txt(s, Inches(0.85), Inches(0.95), Inches(11.5), Inches(0.8),
    [[("数据库设计", 30, FG_STRONG, True, FONT)]])
txt(s, Inches(0.87), Inches(1.67), Inches(11.5), Inches(0.4),
    [[("ORM 定义模型，开发用 SQLite，论文/上线切 PostgreSQL——一份代码两种后端", 14, FG_DIM, False, FONT)]])

# ── Left: simple E-R diagram ──────────────────────────────────────
NW = Inches(2.35)
# column x positions
c1, c2, c3 = Inches(0.6), Inches(3.35), Inches(6.1)
n_user  = node(s, c1, Inches(3.7), NW, "User", "用户（Django 内置）", FG_DIM)
n_stock = node(s, c2, Inches(2.5), NW, "Stock", "证券 code/name", DOWN)
n_bar   = node(s, c2, Inches(4.9), NW, "DailyBar", "日线 OHLCV", DOWN)
n_watch = node(s, c1, Inches(5.3), NW, "WatchItem", "自选股", AI)
n_strat = node(s, c3, Inches(2.5), NW, "Strategy", "策略 + 参数", BRAND)
n_bt    = node(s, c3, Inches(3.9), NW, "Backtest", "回测结果", BRAND)
n_conv  = node(s, c3, Inches(5.3), NW, "Conversation", "对话", UP)
n_msg   = node(s, c3, Inches(6.5), NW, "Message", "消息+trace", UP, h=Inches(0.7))

connect(s, n_stock, n_bar, "1:N", side="down")
connect(s, n_conv, n_msg, "1:N", side="down")
connect(s, n_user, n_watch, "1:N", side="down")

# ── Right: PostgreSQL usage panel ─────────────────────────────────
px = Inches(8.75)
rect(s, px, Inches(2.5), Inches(4.1), Inches(4.5), fill=PANEL, line=BRAND, line_w=1.5)
txt(s, px + Inches(0.3), Inches(2.7), Inches(3.6), Inches(0.5),
    [[("如何用 PostgreSQL", 16, BRAND, True, FONT)]])
pg = [("ORM 屏蔽差异",
       "Django ORM 定义模型，换库只改\nsettings 的 DATABASES，模型代码不动"),
      ("开发→生产平滑迁移",
       "开发用 SQLite 零配置，论文/上线\n用 PostgreSQL，migrate 一键建表"),
      ("JSON 字段更适合 PG",
       "metrics / equity_curve / trace 用\nJSONField，PG 原生 jsonb 可索引查询"),
      ("并发与规模",
       "PG 支持多连接并发写入，\n行情数据量增长后仍稳定")]
py = Inches(3.3)
for i, (h, b) in enumerate(pg):
    y = py + Emu(int(Inches(0.92).emu * i))
    txt(s, px + Inches(0.3), y, Inches(3.6), Inches(0.35),
        [[("● ", 12, DOWN, True, FONT), (h, 13, FG_STRONG, True, FONT)]])
    txt(s, px + Inches(0.52), y + Inches(0.34), Inches(3.45), Inches(0.6),
        [[(ln, 10.5, FG_DIM, False, FONT)] for ln in b.split("\n")], spacing=1.1)

import os
out = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "docs", "数据库设计.pptx"))
prs.save(out)
print("SAVED_OK slides=1")


