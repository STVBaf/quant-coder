"""Generate the defense presentation deck (docs/答辩PPT.pptx).

Dark OKX-style financial theme, 11 dense reading-first slides. Self-contained:
run once with `python gen_deck.py` in the llmbase env.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Palette (mirrors frontend main.css) ──────────────────────────
BG        = RGBColor(0x0B, 0x0E, 0x11)
PANEL     = RGBColor(0x18, 0x1A, 0x20)
LINE      = RGBColor(0x2B, 0x31, 0x39)
FG        = RGBColor(0xEA, 0xEC, 0xEF)
FG_STRONG = RGBColor(0xFA, 0xFA, 0xFA)
FG_DIM    = RGBColor(0x84, 0x8E, 0x9C)
UP        = RGBColor(0xF6, 0x46, 0x5D)  # 红涨
DOWN      = RGBColor(0x0E, 0xCB, 0x81)  # 绿跌
BRAND     = RGBColor(0x00, 0x7A, 0xFF)
AI        = RGBColor(0x8A, 0x2B, 0xE2)

FONT      = "Microsoft YaHei"
FONT_MONO = "Consolas"

# 16:9 at 13.333 x 7.5 inches
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def _no_shadow(shape):
    shape.shadow.inherit = False


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        spacing=1.0):
    """Add a textbox. `runs` = list of paragraphs; each paragraph is a list of
    (text, size, color, bold, font) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for (t, size, color, bold, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = fnt
    return tb


def rect(s, x, y, w, h, fill=PANEL, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    _no_shadow(sp)
    return sp


def accent_bar(s, x, y, h, color=BRAND, w=Inches(0.08)):
    """A thin vertical accent bar, used beside section titles."""
    return rect(s, x, y, w, h, fill=color, line=None, shape=MSO_SHAPE.RECTANGLE)


def chrome(s, tag, page):
    """Per-slide chrome: top-left section tag + bottom-right page number."""
    txt(s, Inches(0.6), Inches(0.34), Inches(6), Inches(0.4),
        [[("QuantDesk", 12, BRAND, True, FONT),
          ("  ·  " + tag, 12, FG_DIM, False, FONT)]])
    txt(s, Inches(11.7), Inches(6.95), Inches(1.1), Inches(0.4),
        [[(f"{page:02d} / 11", 11, FG_DIM, False, FONT_MONO)]],
        align=PP_ALIGN.RIGHT)


def title(s, main, sub=None, y=Inches(0.95)):
    accent_bar(s, Inches(0.6), y + Inches(0.05), Inches(0.62))
    txt(s, Inches(0.85), y, Inches(11.5), Inches(0.8),
        [[(main, 30, FG_STRONG, True, FONT)]])
    if sub:
        txt(s, Inches(0.87), y + Inches(0.72), Inches(11.5), Inches(0.4),
            [[(sub, 14, FG_DIM, False, FONT)]])


def card(s, x, y, w, h, head, lines, accent=BRAND, head_size=15, body_size=12):
    """Panel card: colored heading + bullet-ish body lines (list of strings)."""
    rect(s, x, y, w, h)
    accent_bar(s, x + Inches(0.18), y + Inches(0.22), Inches(0.32), color=accent, w=Inches(0.06))
    txt(s, x + Inches(0.36), y + Inches(0.16), w - Inches(0.5), Inches(0.45),
        [[(head, head_size, FG_STRONG, True, FONT)]])
    body = [[(ln, body_size, FG, False, FONT)] for ln in lines]
    txt(s, x + Inches(0.36), y + Inches(0.68), w - Inches(0.6), h - Inches(0.8),
        body, spacing=1.25)


def chip(s, x, y, w, text, color):
    """Small pill label."""
    sp = rect(s, x, y, w, Inches(0.42), fill=None, line=color, line_w=1.25)
    txt(s, x, y + Inches(0.02), w, Inches(0.38),
        [[(text, 12, color, True, FONT)]], align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    return sp


# ══════════════════════════════════════════════════════════════════
# Slide 1 — Cover
# ══════════════════════════════════════════════════════════════════
s = slide()
# brand glow accents (soft rounded rects behind title)
g = rect(s, Inches(-1.5), Inches(-1.2), Inches(6), Inches(4), fill=RGBColor(0x0E,0x1A,0x2E), line=None, shape=MSO_SHAPE.OVAL)
g2 = rect(s, Inches(9), Inches(4.5), Inches(6), Inches(4), fill=RGBColor(0x1A,0x10,0x28), line=None, shape=MSO_SHAPE.OVAL)
txt(s, Inches(1.2), Inches(1.55), Inches(2), Inches(1),
    [[("Q", 40, FG_STRONG, True, FONT)]])
rect(s, Inches(1.2), Inches(1.5), Inches(0.95), Inches(0.95), fill=None, line=BRAND, line_w=2)
txt(s, Inches(1.15), Inches(2.85), Inches(11), Inches(1.3),
    [[("股票量化研究台", 52, FG_STRONG, True, FONT)]])
txt(s, Inches(1.2), Inches(4.15), Inches(11), Inches(0.6),
    [[("看大盘 → 选股 → 策略回测 → AI Agent 智能解读", 20, BRAND, False, FONT)]])
txt(s, Inches(1.2), Inches(4.85), Inches(11), Inches(0.5),
    [[("个人量化研究台 · 前后端分离 · LLM 驱动回测", 15, FG_DIM, False, FONT)]])
rect(s, Inches(1.22), Inches(5.7), Inches(3.4), Inches(0.02), fill=LINE, line=None, shape=MSO_SHAPE.RECTANGLE)
txt(s, Inches(1.2), Inches(5.9), Inches(11), Inches(0.5),
    [[("汇报人：王颢竣", 16, FG, False, FONT), ("     课程结课项目答辩", 16, FG_DIM, False, FONT)]])


# ══════════════════════════════════════════════════════════════════
# Slide 2 — Background & Positioning
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "选题背景与定位", 2)
title(s, "为什么做这个系统", "个人量化研究台：把专业量化流程搬进一个轻量 Web 应用")
# left: pain points
card(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(2.0), "痛点",
     ["· 量化研究工具专业、门槛高，个人难上手",
      "· 免费行情数据源慢、不稳，采集是第一道坎",
      "· 回测容易踩“未来函数”陷阱，结果不可信",
      "· 策略调参与结果解读，对新手不友好"], accent=UP)
card(s, Inches(0.6), Inches(4.3), Inches(5.9), Inches(2.1), "定位",
     ["· 聚焦研究 / 回测 / 智能分析，不做实盘交易",
      "· 前后端分离，暗色金融风，交互直观",
      "· 用 AI Agent 降低“提出回测需求”的门槛",
      "· 课程结课项目，兼顾工程完整度与论文深度"], accent=DOWN)
# right: flow pipeline
rect(s, Inches(6.9), Inches(2.1), Inches(5.85), Inches(4.3))
txt(s, Inches(7.2), Inches(2.35), Inches(5), Inches(0.5),
    [[("核心流程", 15, FG_STRONG, True, FONT)]])
flow = [("看大盘", "实时指数 / 涨跌广度 / 板块热力图", BRAND),
        ("选股", "搜索个股 · K线 · 加入自选", FG),
        ("策略回测", "三种策略 · 向量化引擎 · 绩效指标", DOWN),
        ("AI 解读", "自然语言驱动回测 · 结果解读", AI)]
fy = Inches(2.95)
for i, (name, desc, col) in enumerate(flow):
    y = fy + Emu(int(Inches(0.86).emu * i))
    rect(s, Inches(7.2), y, Inches(5.25), Inches(0.72), fill=BG, line=col, line_w=1.25)
    txt(s, Inches(7.45), y + Inches(0.06), Inches(2), Inches(0.6),
        [[(name, 16, col, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(9.2), y + Inches(0.06), Inches(3.1), Inches(0.6),
        [[(desc, 11, FG_DIM, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        txt(s, Inches(9.55), y + Inches(0.6), Inches(0.5), Inches(0.3),
            [[("▼", 11, col, False, FONT)]], align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# Slide 3 — Tech Stack
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "技术栈", 3)
title(s, "技术选型", "前后端分离 + 免费数据源 + LLM，围绕“研究台”定位取舍")
stack = [("后端", "Django 6 + DRF", "JSON API · conda llmbase", BRAND),
         ("前端", "Vue 3 + Vite", "Pinia · ECharts · Tailwind v4", DOWN),
         ("数据库", "SQLite / PostgreSQL", "行情/策略/回测/对话持久化", FG),
         ("数据源", "通达信 + akshare", "多源降级 · DB 缓存", UP),
         ("回测", "自写 pandas 向量化", "非调库 · 防未来函数", DOWN),
         ("AI Agent", "Claude + tool use", "自然语言驱动回测", AI)]
cw, ch = Inches(3.9), Inches(1.75)
gx, gy = Inches(0.6), Inches(2.15)
for i, (tag, tech, desc, col) in enumerate(stack):
    r, c = divmod(i, 3)
    x = gx + Emu(int((cw.emu + Inches(0.13).emu) * c))
    y = gy + Emu(int((ch.emu + Inches(0.2).emu) * r))
    rect(s, x, y, cw, ch)
    accent_bar(s, x + Inches(0.2), y + Inches(0.24), Inches(1.25), color=col, w=Inches(0.06))
    txt(s, x + Inches(0.38), y + Inches(0.2), cw - Inches(0.5), Inches(0.4),
        [[(tag, 13, col, True, FONT)]])
    txt(s, x + Inches(0.38), y + Inches(0.62), cw - Inches(0.5), Inches(0.5),
        [[(tech, 17, FG_STRONG, True, FONT)]])
    txt(s, x + Inches(0.38), y + Inches(1.15), cw - Inches(0.5), Inches(0.5),
        [[(desc, 11, FG_DIM, False, FONT)]])


# ══════════════════════════════════════════════════════════════════
# Slide 4 — System Architecture
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "系统架构", 4)
title(s, "前后端分离架构", "四个后端 app 各司其职，行情落库缓存，Key 只在后端")
# frontend box
rect(s, Inches(3.4), Inches(2.05), Inches(6.5), Inches(0.75), fill=BG, line=DOWN, line_w=1.5)
txt(s, Inches(3.4), Inches(2.1), Inches(6.5), Inches(0.65),
    [[("Vue 3 前端", 16, DOWN, True, FONT),
      ("   看板 / 工作台 / 回测历史 / 个股 / 登录", 11, FG_DIM, False, FONT)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(3.4), Inches(2.82), Inches(6.5), Inches(0.4),
    [[("▼   HTTP / JSON · Vite 代理 /api", 11, FG_DIM, False, FONT_MONO)]],
    align=PP_ALIGN.CENTER)
# DRF box
rect(s, Inches(3.4), Inches(3.25), Inches(6.5), Inches(0.7), fill=BG, line=BRAND, line_w=1.5)
txt(s, Inches(3.4), Inches(3.28), Inches(6.5), Inches(0.62),
    [[("Django + DRF 后端", 16, BRAND, True, FONT)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# four services
svc = [("market\n行情服务", "多源降级\n+DB缓存", BRAND),
       ("backtest\n回测引擎", "pandas\n向量化", DOWN),
       ("agent\nAI 服务", "Claude\ntool use", AI),
       ("accounts\n用户系统", "Token 认证\n自选/历史", FG)]
sw = Inches(2.9)
sx0 = Inches(0.85)
for i, (name, desc, col) in enumerate(svc):
    x = sx0 + Emu(int((sw.emu + Inches(0.15).emu) * i))
    rect(s, x, Inches(4.35), sw, Inches(1.3), line=col, line_w=1.25)
    nm = name.split("\n")
    txt(s, x + Inches(0.1), Inches(4.5), sw - Inches(0.2), Inches(0.7),
        [[(nm[0], 12, col, True, FONT_MONO)], [(nm[1], 14, FG_STRONG, True, FONT)]],
        align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(5.15), sw - Inches(0.2), Inches(0.5),
        [[(l, 10.5, FG_DIM, False, FONT)] for l in desc.split("\n")],
        align=PP_ALIGN.CENTER)
txt(s, Inches(0.85), Inches(5.75), Inches(11.6), Inches(0.4),
    [[("▼", 12, FG_DIM, False, FONT)]], align=PP_ALIGN.CENTER)
# DB
rect(s, Inches(3.9), Inches(6.15), Inches(5.5), Inches(0.72), fill=BG, line=LINE, line_w=1.5)
txt(s, Inches(3.9), Inches(6.18), Inches(5.5), Inches(0.64),
    [[("SQLite / PostgreSQL", 14, FG, True, FONT_MONO),
      ("   行情缓存·策略·回测·对话", 11, FG_DIM, False, FONT)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════
# Slide 5 — Module 1: Market Data Layer
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "模块① 行情数据层", 5)
title(s, "多源降级 + 缓存（稳定性设计）", "按数据类型分主备，而非盲目堆一个源——这是工程深度所在")
# table header
tx = Inches(0.6)
cols = [Inches(2.4), Inches(2.9), Inches(2.6), Inches(4.35)]
heads = ["数据类型", "主源", "备源", "说明"]
hx = tx
rect(s, tx, Inches(2.1), Inches(12.25), Inches(0.5), fill=PANEL, line=None)
for hd, cwd in zip(heads, cols):
    txt(s, hx + Inches(0.15), Inches(2.14), cwd, Inches(0.42),
        [[(hd, 13, BRAND, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    hx = hx + cwd
rows = [("指数实时行情", "通达信 mootdx", "akshare 日线", "TDX 原生协议直连券商服务器，免 key、实时、稳"),
        ("个股实时报价", "通达信 mootdx", "—", "给看板 ticker 跑马灯用"),
        ("涨跌广度/板块", "akshare eastmoney", "2 分钟缓存", "数据最全（领涨股/换手率），失败保留上次"),
        ("个股 K 线历史", "akshare 新浪源", "—", "首次拉取落库 SQLite，命中缓存不再打外部")]
ry = Inches(2.62)
for i, row in enumerate(rows):
    rect(s, tx, ry, Inches(12.25), Inches(0.62), fill=(BG if i % 2 else PANEL), line=None)
    cx = tx
    for j, (val, cwd) in enumerate(zip(row, cols)):
        col = FG_STRONG if j == 0 else (BRAND if j == 1 and val != "—" else FG)
        txt(s, cx + Inches(0.15), ry + Inches(0.04), cwd, Inches(0.55),
            [[(val, 11.5, col, j <= 1, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        cx = cx + cwd
    ry = ry + Inches(0.62)
# bottom detail cards
card(s, Inches(0.6), Inches(5.35), Inches(3.95), Inches(1.6), "缓存策略",
     ["· K线落库，覆盖区间才命中", "· 板块 2 分钟 TTL", "· 失败返回 last-good，不空屏"], accent=DOWN, body_size=11.5)
card(s, Inches(4.75), Inches(5.35), Inches(3.95), Inches(1.6), "容错设计",
     ["· 外部调用 3 次退避重试", "· 看板逐源取数", "· 单源失败记 errors 不 500"], accent=UP, body_size=11.5)
card(s, Inches(8.9), Inches(5.35), Inches(3.95), Inches(1.6), "论文对应",
     ["· “行情数据采集与", "  缓存层设计” 章节", "· 可讲工程取舍与韧性"], accent=AI, body_size=11.5)


# ══════════════════════════════════════════════════════════════════
# Slide 6 — Module 2: Backtest Engine
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "模块② 回测引擎", 6)
title(s, "向量化回测引擎", "自写而非调库，shift 防未来函数——回测可信度的关键")
# left: anti-lookahead illustration
rect(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(2.55))
txt(s, Inches(0.85), Inches(2.28), Inches(5.5), Inches(0.4),
    [[("防未来函数（look-ahead bias）", 15, DOWN, True, FONT)]])
txt(s, Inches(0.85), Inches(2.78), Inches(5.5), Inches(0.8),
    [[("信号在 T 日产生，持仓 shift(1) 推迟一天兑现：", 12, FG, False, FONT)],
     [("“今天的信号、明天才成交”", 13, FG_STRONG, True, FONT)]], spacing=1.2)
txt(s, Inches(0.85), Inches(3.72), Inches(5.5), Inches(0.5),
    [[("strat_ret = position.shift(1) * daily_ret", 13, BRAND, False, FONT_MONO)]])
txt(s, Inches(0.85), Inches(4.15), Inches(5.5), Inches(0.4),
    [[("杜绝用到未来信息，结果才可信", 11.5, FG_DIM, False, FONT)]])
# left-bottom: strategies
card(s, Inches(0.6), Inches(4.8), Inches(6.0), Inches(2.05), "三种经典策略",
     ["· 双均线 ma_cross：快线上穿慢线持有",
      "· RSI：超卖买入 / 超买卖出（均值回归）",
      "· 布林带 bollinger：触下轨买入 / 回中轨卖出"], accent=BRAND, body_size=12)
# right: four metrics
txt(s, Inches(6.95), Inches(2.1), Inches(6), Inches(0.4),
    [[("四大绩效指标", 15, DOWN, True, FONT)]])
metrics = [("累计收益", "策略最终净值 − 1", DOWN),
           ("最大回撤", "净值从峰值最大跌幅 · 衡量风险", UP),
           ("夏普比率", "年化超额收益 / 波动率 (×√252)", BRAND),
           ("胜率", "盈利交易日占比", FG)]
my = Inches(2.6)
for i, (name, desc, col) in enumerate(metrics):
    r, c = divmod(i, 2)
    x = Inches(6.95) + Emu(int(Inches(3.05).emu * c))
    y = my + Emu(int(Inches(1.05).emu * r))
    rect(s, x, y, Inches(2.85), Inches(0.9), line=col, line_w=1.25)
    txt(s, x + Inches(0.2), y + Inches(0.1), Inches(2.5), Inches(0.4),
        [[(name, 14, col, True, FONT)]])
    txt(s, x + Inches(0.2), y + Inches(0.48), Inches(2.55), Inches(0.4),
        [[(desc, 9.5, FG_DIM, False, FONT)]])
# right-bottom result callout
rect(s, Inches(6.95), Inches(4.8), Inches(5.9), Inches(2.05), fill=PANEL, line=DOWN, line_w=1.5)
txt(s, Inches(7.2), Inches(5.0), Inches(5.4), Inches(0.4),
    [[("实测样例", 13, DOWN, True, FONT)]])
txt(s, Inches(7.2), Inches(5.45), Inches(5.4), Inches(1.2),
    [[("茅台 · 布林带 · 2023–2024", 15, FG_STRONG, True, FONT)],
     [("累计收益 ", 13, FG, False, FONT), ("+21.6%", 20, UP, True, FONT_MONO),
      ("    夏普 ", 13, FG, False, FONT), ("0.78", 20, DOWN, True, FONT_MONO)]],
    spacing=1.3)


# ══════════════════════════════════════════════════════════════════
# Slide 7 — Module 3: AI Agent (HIGHLIGHT)
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "模块③ 量化 AI Agent · 核心亮点", 7)
# highlight accent uses AI purple
accent_bar(s, Inches(0.6), Inches(1.0), Inches(0.62), color=AI)
txt(s, Inches(0.85), Inches(0.95), Inches(11.5), Inches(0.8),
    [[("量化 AI Agent", 30, FG_STRONG, True, FONT),
      ("   核心亮点", 16, AI, True, FONT)]])
txt(s, Inches(0.87), Inches(1.67), Inches(11.5), Inches(0.4),
    [[("Claude + tool use 手动循环：把大模型真正“用起来”，而非套个聊天框", 14, FG_DIM, False, FONT)]])
# tool-use loop
rect(s, Inches(0.6), Inches(2.3), Inches(7.4), Inches(2.5), line=AI, line_w=1.5)
txt(s, Inches(0.85), Inches(2.48), Inches(6.8), Inches(0.4),
    [[("tool use 手动循环", 15, AI, True, FONT)]])
loop = [("① 用户自然语言", "“用双均线回测茅台近三年，5日20日均线”"),
        ("② Claude 解析", "拆出 code / kind / params，决定调哪个工具"),
        ("③ 后端执行工具", "run_backtest 复用向量化回测引擎"),
        ("④ 结果喂回 Claude", "tool_result 回传，循环至 end_turn"),
        ("⑤ 中文解读", "收益/回撤/夏普/风险，数字全部来自工具")]
ly = Inches(2.92)
for i, (a, b) in enumerate(loop):
    y = ly + Emu(int(Inches(0.37).emu * i))
    txt(s, Inches(0.85), y, Inches(2.55), Inches(0.35),
        [[(a, 12, FG_STRONG, True, FONT)]])
    txt(s, Inches(3.45), y, Inches(4.4), Inches(0.35),
        [[(b, 10.5, FG_DIM, False, FONT)]])
# right: tools + design
card(s, Inches(8.25), Inches(2.3), Inches(4.55), Inches(1.35), "三个工具（复用后端能力）",
     ["· list_strategies 列策略",
      "· get_quote 查行情",
      "· run_backtest 跑回测"], accent=BRAND, body_size=11.5)
card(s, Inches(8.25), Inches(3.78), Inches(4.55), Inches(1.02), "工程要点",
     ["· 对话持久化，前端可展开 trace",
      "· prompt caching 降 token 成本"], accent=DOWN, body_size=11.5)
# bottom band: security + why it matters
rect(s, Inches(0.6), Inches(4.95), Inches(12.25), Inches(1.9), fill=PANEL, line=None)
txt(s, Inches(0.85), Inches(5.12), Inches(12), Inches(0.4),
    [[("为什么是亮点", 14, AI, True, FONT)]])
txt(s, Inches(0.85), Inches(5.55), Inches(12), Inches(1.2),
    [[("· 数字不编造：system prompt 强约束所有数字必须来自工具返回，杜绝幻觉", 12.5, FG, False, FONT)],
     [("· Key 只在后端 .env，绝不下发前端——安全边界清晰", 12.5, FG, False, FONT)],
     [("· 区别于普通课设：完整跑通“自然语言 → 工具调用 → 结果解读”闭环", 12.5, FG_STRONG, True, FONT)]],
    spacing=1.35)


# ══════════════════════════════════════════════════════════════════
# Slide 8 — Module 4: Dashboard UI
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "模块④ 看板 UI", 8)
title(s, "极客风大盘看板", "OKX 风暗色金融界面 + 实时轮询 + 动效，演示加分项")
card(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(2.15), "实时与交互",
     ["· 前端定时 8s 轮询大盘快照与报价",
      "· 数字平滑 count-up + 红涨绿跌 tick 闪烁",
      "· 顶部 ticker 跑马灯，热门股实时滚动",
      "· 涨跌广度做成“恐贪指数”半圆仪表盘"], accent=BRAND, body_size=12.5)
card(s, Inches(0.6), Inches(4.4), Inches(6.0), Inches(2.0), "视觉设计",
     ["· A 股惯例红涨绿跌，等宽字体显示数字",
      "· 科技网格背景 + 蓝紫渐变光晕",
      "· 板块热力图 treemap，只展示涨跌前 20",
      "· 错峰进场动画，玻璃拟态面板"], accent=AI, body_size=12.5)
# right: swatches + component list
rect(s, Inches(6.9), Inches(2.1), Inches(5.95), Inches(4.3))
txt(s, Inches(7.2), Inches(2.3), Inches(5), Inches(0.4),
    [[("设计语言 · 配色", 15, FG_STRONG, True, FONT)]])
sw_data = [("背景 #0B0E11", BG), ("面板 #181A20", PANEL),
           ("涨 #F6465D", UP), ("跌 #0ECB81", DOWN),
           ("品牌 #007AFF", BRAND), ("AI #8A2BE2", AI)]
for i, (lbl, col) in enumerate(sw_data):
    r, c = divmod(i, 2)
    x = Inches(7.2) + Emu(int(Inches(2.85).emu * c))
    y = Inches(2.85) + Emu(int(Inches(0.7).emu * r))
    rect(s, x, y, Inches(0.5), Inches(0.5), fill=col, line=LINE, shape=MSO_SHAPE.RECTANGLE)
    txt(s, x + Inches(0.62), y, Inches(2.1), Inches(0.5),
        [[(lbl, 12, FG, False, FONT_MONO)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(7.2), Inches(5.15), Inches(5.3), Inches(0.4),
    [[("五个前端视图", 14, BRAND, True, FONT)]])
txt(s, Inches(7.2), Inches(5.6), Inches(5.4), Inches(1.1),
    [[("大盘看板 · 量化工作台 · 回测历史", 12.5, FG, False, FONT)],
     [("个股详情 · 登录注册", 12.5, FG, False, FONT)]], spacing=1.3)


# ══════════════════════════════════════════════════════════════════
# Slide 9 — Difficulties & Solutions
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "技术难点与解决", 9)
title(s, "技术难点与解决", "答辩追问预案：每个难点都有明确的工程解法")
diff = [("免费数据源不稳、慢", "多源降级 + DB 缓存 + 退避重试 + 逐源容错", BRAND),
        ("回测的未来函数陷阱", "持仓 shift(1)，信号 T 日产生、T+1 兑现", DOWN),
        ("让 LLM 真正“做事”", "tool use 手动循环 + 工具复用引擎 + 数字强约束", AI),
        ("通达信无行业板块涨跌幅", "技术验证后务实取舍：行业板块保留 akshare", UP),
        ("看板实时感", "前端轮询 + 数字动画，而非整页刷新", FG)]
dy = Inches(2.15)
for i, (problem, solution, col) in enumerate(diff):
    y = dy + Emu(int(Inches(0.92).emu * i))
    rect(s, Inches(0.6), y, Inches(12.25), Inches(0.78), fill=(PANEL if i % 2 == 0 else BG), line=None)
    accent_bar(s, Inches(0.6), y, Inches(0.78), color=col, w=Inches(0.08))
    txt(s, Inches(0.95), y + Inches(0.04), Inches(4.4), Inches(0.7),
        [[("难点  ", 10, FG_DIM, False, FONT), (problem, 14, FG_STRONG, True, FONT)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.6), y + Inches(0.04), Inches(0.6), Inches(0.7),
        [[("→", 15, col, True, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.2), y + Inches(0.04), Inches(6.5), Inches(0.7),
        [[(solution, 12.5, FG, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════
# Slide 10 — Results & Live Demo
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "实验结果与演示", 10)
title(s, "回测实验结果 & Live Demo", "真实 A 股数据跑通，现场演示四步走")
# left: results
rect(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.75))
txt(s, Inches(0.85), Inches(2.3), Inches(5), Inches(0.4),
    [[("回测实验（真实行情）", 15, DOWN, True, FONT)]])
res = [("茅台 · 布林带 · 23–24", "+21.6%", "夏普 0.78", UP),
       ("向量化引擎", "一次算完整曲线", "非逐日循环", BRAND),
       ("防未来函数", "shift(1) 已验证", "结果可信", DOWN)]
for i, (name, v1, v2, col) in enumerate(res):
    y = Inches(2.8) + Emu(int(Inches(1.25).emu * i))
    rect(s, Inches(0.85), y, Inches(5.5), Inches(1.05), fill=BG, line=col, line_w=1.25)
    txt(s, Inches(1.1), y + Inches(0.12), Inches(5), Inches(0.4),
        [[(name, 13, FG_STRONG, True, FONT)]])
    txt(s, Inches(1.1), y + Inches(0.52), Inches(5), Inches(0.45),
        [[(v1, 15, col, True, FONT_MONO), ("   " + v2, 12, FG_DIM, False, FONT)]])
# right: demo steps
rect(s, Inches(6.9), Inches(2.1), Inches(5.95), Inches(4.75), fill=PANEL, line=BRAND, line_w=1.25)
txt(s, Inches(7.2), Inches(2.3), Inches(5), Inches(0.4),
    [[("现场演示脚本", 15, BRAND, True, FONT)]])
demo = [("1  大盘看板", "指数跳动 · ticker 滚动 · 热力图，视觉开场"),
        ("2  个股回测", "搜茅台 → 布林带 → 收益曲线 + 买卖点"),
        ("3  AI Agent", "自然语言 → 自动回测 → 解读 → 展开 trace"),
        ("4  用户系统", "登录后回测入历史，历史页对比多次结果")]
for i, (a, b) in enumerate(demo):
    y = Inches(2.85) + Emu(int(Inches(0.98).emu * i))
    col = AI if i == 2 else FG_STRONG
    txt(s, Inches(7.25), y, Inches(5.3), Inches(0.4),
        [[(a, 15, col, True, FONT), ("   ← 高潮" if i == 2 else "", 11, AI, True, FONT)]])
    txt(s, Inches(7.25), y + Inches(0.4), Inches(5.35), Inches(0.5),
        [[(b, 11, FG_DIM, False, FONT)]])


# ══════════════════════════════════════════════════════════════════
# Slide 11 — Summary & Outlook
# ══════════════════════════════════════════════════════════════════
s = slide()
chrome(s, "总结与展望", 11)
title(s, "总结与展望", "五个里程碑全部交付，后续可继续深化")
# delivered stats
stats = [("4", "后端 app"), ("11", "API 接口"), ("3", "回测策略"),
         ("3", "Agent 工具"), ("5", "前端视图")]
for i, (num, lbl) in enumerate(stats):
    x = Inches(0.6) + Emu(int(Inches(2.5).emu * i))
    rect(s, x, Inches(2.15), Inches(2.3), Inches(1.35))
    txt(s, x, Inches(2.35), Inches(2.3), Inches(0.7),
        [[(num, 34, BRAND, True, FONT_MONO)]], align=PP_ALIGN.CENTER)
    txt(s, x, Inches(3.02), Inches(2.3), Inches(0.4),
        [[(lbl, 12, FG_DIM, False, FONT)]], align=PP_ALIGN.CENTER)
card(s, Inches(0.6), Inches(3.85), Inches(6.0), Inches(2.35), "已交付（M1–M5）",
     ["· M1 脚手架与数据层", "· M2 大盘看板", "· M3 回测引擎",
      "· M4 用户系统与历史", "· M5 量化 Agent"], accent=DOWN, body_size=12.5)
card(s, Inches(6.85), Inches(3.85), Inches(6.0), Inches(2.35), "未来展望",
     ["· Agent 流式输出（打字机效果）", "· 板块迁宽基指数 / 更多数据源",
      "· 更多策略与多股组合回测", "· PostgreSQL + 部署上线",
      "· 策略参数自动寻优"], accent=AI, body_size=12.5)
txt(s, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.6),
    [[("谢谢聆听 · 敬请指正", 20, FG_STRONG, True, FONT)]], align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(__file__), "..", "docs", "答辩PPT.pptx")
out = os.path.abspath(out)
prs.save(out)
print("SAVED", out, "slides=", len(prs.slides._sldIdLst))







